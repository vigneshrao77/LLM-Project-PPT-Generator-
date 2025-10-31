import os
import re
import requests
import urllib.parse
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from ollama import chat
from PIL import Image
import streamlit as st
import time
import json
import textwrap

# ======================================================
# UI / Page Style 
st.set_page_config(page_title="Paper2PPT", page_icon="🎓", layout="centered")

custom_css = """
<style>
body {
  background: radial-gradient(circle at 50% 40%, #021325 0%, #022f46 50%, #001f3f 100%);
  background-size: 300% 300%;
  animation: gradientMove 15s ease infinite;
  font-family: 'Poppins', sans-serif;
  color: white;
}
@keyframes gradientMove {
  0% { background-position: 0% 50%; }
  50% { background-position: 100% 50%; }
  100% { background-position: 0% 50%; }
}
h1 { text-align: center; color: #9fb8ff; font-weight: 800; font-size: 2.4rem; margin-top: -18px; text-shadow: 0 0 18px rgba(159,184,255,0.18); }
h2, h3, p { text-align: center; color: #dfefff; }
.stTextInput>div>div>input, .stTextArea textarea {
  background: rgba(255, 255, 255, 0.04);
  border: 1px solid rgba(159,184,255,0.08);
  border-radius: 14px;
  color: white;
  font-size: 1rem;
  transition: all 0.25s ease;
  padding: 0.6rem;
}
.stTextInput>div>div>input:hover, .stTextArea textarea:hover {
  border: 1px solid rgba(159,184,255,0.22);
  box-shadow: 0 0 18px rgba(159,184,255,0.06);
}
.stButton>button {
  background: linear-gradient(135deg, #4b6ef6, #6a9bff);
  color: white;
  border: none;
  border-radius: 12px;
  padding: 0.6rem 1.3rem;
  font-weight: 600;
  font-size: 1rem;
  cursor: pointer;
  transition: all 0.25s ease;
  box-shadow: 0px 0px 18px rgba(75,110,246,0.18);
}
.stButton>button:hover {
  transform: scale(1.04);
  box-shadow: 0px 6px 26px rgba(75,110,246,0.28);
}
div[data-testid="stSpinner"] { color: #bfe0ff !important; text-align:center !important; }
.small-glow { color: #cfe8ff; text-align:center; font-size:0.9rem; text-shadow: 0 0 8px rgba(159,184,255,0.08); }
.success-box { background: rgba(46, 125, 50, 0.2); border: 1px solid #2e7d32; border-radius: 10px; padding: 15px; margin: 10px 0; }
.paper-info { background: rgba(30, 136, 229, 0.15); border: 1px solid #1e88e5; border-radius: 10px; padding: 15px; margin: 10px 0; }
</style>
"""
st.markdown(custom_css, unsafe_allow_html=True)

st.markdown("<h1>🎓 Paper2PPT</h1>", unsafe_allow_html=True)
st.markdown("<p>AI-Enhanced Presentations from research papers using Ollama 3.2</p>", unsafe_allow_html=True)

#  Safety Layer
class SafetyLayer:
    def __init__(self):
        self.bad_words = ["fuck", "shit", "bitch", "nigger", "kill", "hate"]
        self.ignore_phrases = [
            "here are", "below are", "these are", "in this section",
            "bullet points", "slide content", "summary of", "this section covers",
            "here is", "here's"
        ]
    def clean(self, text):
        if not text:
            return ""
        cleaned = text.strip()
        for phrase in self.ignore_phrases:
            if cleaned.lower().startswith(phrase):
                cleaned = cleaned[len(phrase):].strip()
        for token in ["```", "###", "{", "}", "[", "]"]:
            cleaned = cleaned.replace(token, "")
        for bad in self.bad_words:
            cleaned = cleaned.replace(bad, "*")
        return " ".join(cleaned.split())

safety_layer = SafetyLayer()


#  Enhanced Paper Fetch (Multi-source with fallbacks)
def get_paper_data(paper_url_or_query, ui=None):
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36",
        "Accept": "application/json"
    }
    
    query = paper_url_or_query.strip()
    
    # Progress tracking
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    def update_status(step, total_steps, message):
        progress_bar.progress(step / total_steps)
        status_text.text(f"🔍 {message}")
        time.sleep(0.5)  # Small delay for better UX
    
    # Try multiple strategies
    strategies = [
        ("Semantic Scholar Direct", lambda: try_semantic_scholar_direct(query, headers)),
        ("Semantic Scholar Search", lambda: try_semantic_scholar_search(query, headers)),
        ("CrossRef Search", lambda: try_crossref_search(query, headers)),
        ("arXiv Direct", lambda: try_arxiv_direct(query, headers)),
        ("arXiv Search", lambda: try_arxiv_search(query, headers)),
        ("Google Scholar Simulation", lambda: try_google_scholar_simulation(query, headers))
    ]
    
    for i, (strategy_name, strategy_func) in enumerate(strategies):
        update_status(i, len(strategies), f"Trying {strategy_name}...")
        try:
            result = strategy_func()
            if result and result.get("title") and result.get("abstract"):
                progress_bar.progress(100)
                status_text.text(f"✅ Found via {strategy_name}!")
                st.markdown(f'<div class="success-box">✅ <strong>Paper Found:</strong> "{result["title"]}"</div>', unsafe_allow_html=True)
                return result
        except Exception as e:
            continue
    
    progress_bar.empty()
    status_text.empty()
    return None
#====================================================================================
def try_semantic_scholar_direct(query, headers):
    """Try direct Semantic Scholar paper ID extraction"""
    base = "https://api.semanticscholar.org/graph/v1"
    
    # Extract paper ID from various Semantic Scholar URL formats
    patterns = [
        r"/paper/([0-9a-f]+)",
        r"paperId=([0-9a-f]+)",
        r"SemanticScholarID:([0-9a-f]+)"
    ]
    
    paper_id = None
    for pattern in patterns:
        match = re.search(pattern, query)
        if match:
            paper_id = match.group(1)
            break
    
    if paper_id:
        url = f"{base}/paper/{paper_id}?fields=title,abstract,authors,year,url,venue,publicationVenue,citationCount"
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code == 200:
            data = response.json()
            if data.get("title"):
                return format_semantic_scholar_data(data)
    
    return None

def try_semantic_scholar_search(query, headers):
    """Search Semantic Scholar"""
    base = "https://api.semanticscholar.org/graph/v1"
    
    # Clean query for search
    search_query = re.sub(r'https?://[^\s]+', '', query).strip()
    if not search_query or len(search_query) < 5:
        return None
    
    params = {
        "query": search_query,
        "limit": 5,
        "fields": "title,abstract,authors,year,url,venue,publicationVenue,citationCount"
    }
    
    response = requests.get(f"{base}/paper/search", params=params, headers=headers, timeout=10)
    if response.status_code == 200:
        data = response.json().get("data", [])
        for paper in data:
            if paper.get("title") and paper.get("abstract"):
                return format_semantic_scholar_data(paper)
    
    return None

def try_crossref_search(query, headers):
    """Search CrossRef"""
    search_query = re.sub(r'https?://[^\s]+', '', query).strip()
    if not search_query:
        return None
    
    # Try multiple search strategies
    search_urls = [
        f"https://api.crossref.org/works?query.title={urllib.parse.quote_plus(search_query)}&rows=5",
        f"https://api.crossref.org/works?query={urllib.parse.quote_plus(search_query)}&rows=5",
        f"https://api.crossref.org/works?query.bibliographic={urllib.parse.quote_plus(search_query)}&rows=5"
    ]
    
    for url in search_urls:
        try:
            response = requests.get(url, headers=headers, timeout=10)
            if response.status_code == 200:
                items = response.json().get("message", {}).get("items", [])
                for item in items:
                    if item.get("title") and len(item["title"]) > 0:
                        return format_crossref_data(item)
        except:
            continue
    
    return None

def try_arxiv_direct(query, headers):
    """Extract direct arXiv IDs"""
    patterns = [
        r'arxiv\.org/abs/([0-9]+\.[0-9v]+)',
        r'arxiv\.org/pdf/([0-9]+\.[0-9v]+)',
        r'arxiv\.org/abs/([a-z\-]+\.[0-9v]+)',
        r'arxiv\.org/pdf/([a-z\-]+\.[0-9v]+)',
        r'arXiv:([0-9]+\.[0-9v]+)',
        r'arXiv:([a-z\-]+\.[0-9v]+)'
    ]
    
    for pattern in patterns:
        match = re.search(pattern, query, re.IGNORECASE)
        if match:
            arxiv_id = match.group(1)
            return fetch_arxiv_data(arxiv_id, headers)
    
    return None

def try_arxiv_search(query, headers):
    """Search arXiv"""
    search_query = re.sub(r'https?://[^\s]+', '', query).strip()
    if not search_query or len(search_query) < 5:
        return None
    
    # Clean query for arXiv search
    search_query = re.sub(r'[^a-zA-Z0-9\s]', ' ', search_query)
    search_query = ' '.join(search_query.split()[:10])  # Limit length
    
    url = f"http://export.arxiv.org/api/query?search_query=all:{urllib.parse.quote_plus(search_query)}&start=0&max_results=5"
    
    try:
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code == 200:
            entries = re.findall(r'<entry>(.*?)</entry>', response.text, re.DOTALL)
            for entry in entries:
                title_match = re.search(r'<title>(.*?)</title>', entry, re.DOTALL)
                summary_match = re.search(r'<summary>(.*?)</summary>', entry, re.DOTALL)
                id_match = re.search(r'<id>(.*?)</id>', entry)
                
                if title_match and summary_match:
                    title = re.sub(r'\s+', ' ', title_match.group(1)).strip()
                    abstract = re.sub(r'\s+', ' ', summary_match.group(1)).strip()
                    arxiv_id = id_match.group(1).split('/')[-1] if id_match else "Unknown"
                    
                    return {
                        "title": title,
                        "abstract": abstract,
                        "authors": extract_arxiv_authors(entry),
                        "year": extract_arxiv_year(entry),
                        "venue": "arXiv",
                        "url": f"https://arxiv.org/abs/{arxiv_id}",
                        "source": "arXiv Search"
                    }
    except:
        pass
    
    return None

def try_google_scholar_simulation(query, headers):
    """Simulate Google Scholar search using Semantic Scholar as proxy"""
    search_query = re.sub(r'https?://[^\s]+', '', query).strip()
    if not search_query or len(search_query) < 5:
        return None
    
    # Use Semantic Scholar with broader search
    return try_semantic_scholar_search(search_query, headers)

def format_semantic_scholar_data(data):
    """Format Semantic Scholar data consistently"""
    authors = []
    if data.get("authors"):
        authors = [{"name": author.get("name", "Unknown")} for author in data["authors"]]
    
    venue = data.get("venue") or data.get("publicationVenue", {}).get("name", "Unknown")
    
    return {
        "title": data.get("title", ""),
        "abstract": data.get("abstract", ""),
        "authors": authors,
        "year": data.get("year"),
        "venue": venue,
        "url": data.get("url", ""),
        "citationCount": data.get("citationCount", 0),
        "source": "Semantic Scholar"
    }

def format_crossref_data(item):
    """Format CrossRef data consistently"""
    authors = []
    if item.get("author"):
        for author in item["author"]:
            name = ""
            if author.get("given") and author.get("family"):
                name = f"{author['given']} {author['family']}"
            elif author.get("name"):
                name = author["name"]
            elif author.get("given"):
                name = author["given"]
            elif author.get("family"):
                name = author["family"]
            if name:
                authors.append({"name": name})
    
    title = item.get("title", [""])[0] if item.get("title") else ""
    abstract = item.get("abstract", "")
    if abstract and isinstance(abstract, str) and abstract.startswith("<jats:p>"):
        # Clean JATS XML abstract
        abstract = re.sub(r'<[^>]+>', '', abstract)
    
    return {
        "title": title,
        "abstract": abstract,
        "authors": authors,
        "year": item.get("created", {}).get("date-parts", [[None]])[0][0],
        "venue": item.get("container-title", [""])[0] if item.get("container-title") else "",
        "url": item.get("URL", item.get("url", [""])[0] if item.get("url") else ""),
        "source": "CrossRef"
    }

def fetch_arxiv_data(arxiv_id, headers):
    """Fetch arXiv data by ID"""
    url = f"http://export.arxiv.org/api/query?id_list={arxiv_id}"
    try:
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code == 200:
            entry = response.text
            title_match = re.search(r'<title>(.*?)</title>', entry, re.DOTALL)
            summary_match = re.search(r'<summary>(.*?)</summary>', entry, re.DOTALL)
            
            if title_match and summary_match:
                title = re.sub(r'\s+', ' ', title_match.group(1)).strip()
                abstract = re.sub(r'\s+', ' ', summary_match.group(1)).strip()
                
                return {
                    "title": title,
                    "abstract": abstract,
                    "authors": extract_arxiv_authors(entry),
                    "year": extract_arxiv_year(entry),
                    "venue": "arXiv",
                    "url": f"https://arxiv.org/abs/{arxiv_id}",
                    "source": "arXiv Direct"
                }
    except:
        pass
    
    return None

def extract_arxiv_authors(entry):
    """Extract authors from arXiv entry"""
    authors = []
    author_matches = re.findall(r'<author>\s*<name>([^<]+)</name>\s*</author>', entry)
    for author_name in author_matches:
        authors.append({"name": author_name.strip()})
    return authors

def extract_arxiv_year(entry):
    """Extract year from arXiv entry"""
    published_match = re.search(r'<published>([0-9]{4})', entry)
    if published_match:
        return int(published_match.group(1))
    return None

# Ollama Integration with Smart Text Wrapping
def generate_section_content(section, paper, style_prompt):
    title = paper.get("title", "")
    abstract = paper.get("abstract", "")
    authors = ", ".join(a["name"] for a in paper.get("authors", [])) if paper.get("authors") else ""
    
    prompt = f"""
Generate 3-4 concise bullet points for a slide titled '{section}'.
Paper title: {title}
Authors: {authors}
Abstract: {abstract}
Style: {style_prompt}

IMPORTANT: Each bullet point must be UNDER 12 words and fit on one line in a PowerPoint slide.
Return only the bullet points, no explanations.
"""
    try:
        resp = chat(model="llama3.2", messages=[{"role": "user", "content": prompt}])
        raw = resp["message"]["content"].strip()
        bullets = []
        for l in raw.split("\n"):
            clean = safety_layer.clean(l.lstrip("•*-0123456789. ").strip())
            if 5 <= len(clean) <= 80:  # Reduced max length for better fitting
                bullets.append(clean)
        return bullets[:4]
    except Exception as e:
        # Fallback content that definitely fits
        return [
            f"Key insight from {section.lower()}",
            f"Main finding in this area",
            f"Important contribution highlighted",
            f"Summary of {section.lower()} approach"
        ]

# PPT Builder with 8-9 Words Per Line Wrapping
def build_pptx(paper, style_prompt, output_path="presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colors
    BG_COLOR = RGBColor(10, 25, 70)
    ACCENT = RGBColor(255, 180, 30)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(25, 25, 25)
    LIGHT_BG = RGBColor(250, 250, 252)

    # ========== TITLE SLIDE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_COLOR

    # Smart title wrapping - 8-9 words per line
    title = paper.get("title", "Research Presentation")
    wrapped_title = smart_text_wrap_words(title, max_lines=3, words_per_line=8)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    for i, line in enumerate(wrapped_title):
        p = title_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(36 if i == 0 else 32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(5)

    # Authors with wrapping - 8-9 words per line
    if paper.get("authors"):
        author_names = ", ".join([author["name"] for author in paper["authors"][:6]])
        wrapped_authors = smart_text_wrap_words(author_names, max_lines=2, words_per_line=9)
        
        author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(1))
        author_frame = author_box.text_frame
        author_frame.word_wrap = True
        
        for line in wrapped_authors:
            p = author_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 255)
            p.alignment = PP_ALIGN.CENTER

    # ========== CONTENT SLIDES ==========
    sections = ["Introduction", "Objectives", "Methodology", "Results", "Conclusion"]
    
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = LIGHT_BG

        # Slide title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.add_paragraph()
        p.text = sec
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BG_COLOR
        p.alignment = PP_ALIGN.LEFT

        # Accent bar
        bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT

        # Content area with 8-9 words per line wrapping
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Get and add bullet points
        bullets = generate_section_content(sec, paper, style_prompt)
        
        for i, bullet in enumerate(bullets):
            # Wrap bullets to 8-9 words per line
            wrapped_bullet = smart_bullet_wrap_words(bullet, words_per_line=8)
            
            for j, line in enumerate(wrapped_bullet):
                para = content_frame.add_paragraph()
                if j == 0:
                    prefix = "• " 
                else:
                    prefix = "  "  # Indentation for wrapped lines
                
                para.text = prefix + line
                para.font.size = Pt(24)  # Slightly larger font since we have more space
                para.font.color.rgb = DARK_TEXT
                para.space_after = Pt(10)  # More spacing for better readability
                para.line_spacing = 1.3

        # Add page number
        footer = slide.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.4))
        footer_frame = footer.text_frame
        p_footer = footer_frame.add_paragraph()
        p_footer.text = f"{sections.index(sec) + 1}/{len(sections)}"
        p_footer.font.size = Pt(14)
        p_footer.font.color.rgb = RGBColor(150, 150, 150)
        p_footer.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    return output_path
#===============================================================================
def smart_text_wrap_words(text, max_lines=3, words_per_line=8):
    """Smart text wrapping that ensures 8-9 words per line"""
    if not text:
        return [""]
    
    # Clean the text
    text = ' '.join(text.split())
    words = text.split()
    
    # If text is short enough, return as is
    if len(words) <= words_per_line:
        return [text]
    
    lines = []
    current_line = []
    
    for word in words:
        current_line.append(word)
        
        # If we've reached the word limit for this line, save it
        if len(current_line) >= words_per_line:
            lines.append(' '.join(current_line))
            current_line = []
            
            # Check if we've reached max lines
            if len(lines) >= max_lines:
                break
    
    # Add the last line if we have space
    if current_line and len(lines) < max_lines:
        lines.append(' '.join(current_line))
    
    # Handle truncation with ellipsis if we have more content
    if len(lines) == max_lines and len(words) > len(' '.join(lines).split()):
        last_line = lines[-1]
        if len(last_line.split()) >= words_per_line:
            lines[-1] = ' '.join(last_line.split()[:words_per_line-1]) + " ..."
        else:
            lines[-1] += " ..."
    
    return lines

def smart_bullet_wrap_words(bullet_text, words_per_line=8):
    """Wrap bullet text to 8-9 words per line"""
    if not bullet_text:
        return [""]
    
    # Remove any existing bullet characters
    clean_text = bullet_text.lstrip("•*- ").strip()
    words = clean_text.split()
    
    # If text fits in one line, return as is
    if len(words) <= words_per_line:
        return [clean_text]
    
    # Split into lines of 8-9 words each
    lines = []
    for i in range(0, len(words), words_per_line):
        line = ' '.join(words[i:i + words_per_line])
        lines.append(line)
    
    return lines[:3]  # Maximum 3 lines per bullet


#Streamlit UI

paper_input = st.text_input(
    "Enter paper title, DOI, arXiv URL, Semantic Scholar link, or any paper identifier:",
    placeholder="e.g., 'Attention is All You Need' or 'https://arxiv.org/abs/1706.03762'"
)

style_prompt = st.text_area(
    "Describe slide style:",
    "Academic presentation with clear, concise bullet points. Professional tone.",
    placeholder="Describe how you want the slides to look and feel..."
)

if st.button("✨ Generate Presentation", type="primary"):
    if not paper_input.strip():
        st.warning("Please enter a paper title, URL, or identifier.")
    else:
        with st.spinner("🔍 Searching across multiple academic databases..."):
            paper = get_paper_data(paper_input)
            
            if not paper:
                st.error("""
                ❌ Could not retrieve paper data. Please try:
                - A more specific paper title
                - Direct arXiv URL (e.g., https://arxiv.org/abs/1706.03762)
                - Semantic Scholar link
                - DOI identifier
                - Different wording of the paper title
                """)
            else:
                # Display paper info
                st.markdown(f"""
                <div class="paper-info">
                <h3>📄 Paper Found!</h3>
                <p><strong>Title:</strong> {paper['title']}</p>
                <p><strong>Source:</strong> {paper.get('source', 'Unknown')}</p>
                <p><strong>Year:</strong> {paper.get('year', 'Unknown')}</p>
                <p><strong>Venue:</strong> {paper.get('venue', 'Unknown')}</p>
                </div>
                """, unsafe_allow_html=True)
                
                # Generate presentation
                with st.spinner("📊 Creating presentation slides with perfect formatting..."):
                    output_path = "research_presentation.pptx"
                    try:
                        build_pptx(paper, style_prompt, output_path)
                        st.success("✅ Presentation ready! Download below:")
                        
                        with open(output_path, "rb") as f:
                            st.download_button(
                                "📥 Download PowerPoint Presentation",
                                f,
                                file_name="research_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    except Exception as e:
                        st.error(f"❌ Error generating presentation: {str(e)}")

st.markdown("<p class='small-glow'>✨ Multi-source paper fetching • 8-9 words per line • Professional slides</p>", unsafe_allow_html=True)